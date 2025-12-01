import os
import streamlit as st
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
from groq import Groq
import numpy as np
from dotenv import load_dotenv
from pypdf import PdfReader
import tempfile
from docx import Document as DocxDocument
from pptx import Presentation
import requests
from urllib.parse import urlparse, parse_qs, unquote
from requests.auth import HTTPBasicAuth
from requests_ntlm import HttpNtlmAuth
import msal
import time
from functools import wraps

load_dotenv()

# Page configuration
st.set_page_config(
    page_title="Delivery Management System",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
    <style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        text-align: center;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 2rem;
    }
    .stButton>button {
        width: 100%;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        font-weight: bold;
        border-radius: 10px;
        padding: 0.5rem 1rem;
        border: none;
    }
    .stButton>button:hover {
        background: linear-gradient(90deg, #764ba2 0%, #667eea 100%);
        transform: scale(1.02);
    }
    .result-box {
        padding: 1rem;
        margin: 0.5rem 0;
        background: #f0f2f6;
        border-left: 4px solid #667eea;
        border-radius: 5px;
    }
    .ai-response {
        padding: 1.5rem;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin: 1rem 0;
    }
    </style>
""", unsafe_allow_html=True)

# --------------------------
# Config
# --------------------------
QDRANT_URL = os.getenv("QDRANT_URL")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

COLLECTION_NAME = os.getenv("COLLECTION_NAME", "pdf_chunks")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
MAX_SEARCH_RESULTS = int(os.getenv("MAX_SEARCH_RESULTS", 5))
LLM_MODEL = "openai/gpt-oss-20b"

# --------------------------
# Initialize session state
# --------------------------
if 'embed_model' not in st.session_state:
    st.session_state.embed_model = None
if 'qdrant' not in st.session_state:
    st.session_state.qdrant = None
if 'groq_client' not in st.session_state:
    st.session_state.groq_client = None

# --------------------------
# Embeddings
# --------------------------
@st.cache_resource(show_spinner="Loading embedding model...")
def load_model():
    if st.session_state.embed_model is None:
        from sentence_transformers import SentenceTransformer
        st.session_state.embed_model = SentenceTransformer(EMBEDDING_MODEL, device="cpu")
    return st.session_state.embed_model

def embed(text):
    model = load_model()
    v = model.encode(text, normalize_embeddings=True)
    return np.asarray(v, dtype=np.float32)

def get_vector_size():
    """Get the vector size from the embedding model"""
    model = load_model()
    return model.get_sentence_embedding_dimension()

# --------------------------
# Initialize clients
# --------------------------
def init_clients():
    if st.session_state.qdrant is None:
        try:
            st.session_state.qdrant = QdrantClient(
                url=QDRANT_URL,
                api_key=QDRANT_API_KEY
            )
        except Exception as e:
            st.error(f"Failed to connect to Qdrant: {e}")
            return False
    
    if st.session_state.groq_client is None:
        try:
            st.session_state.groq_client = Groq(api_key=GROQ_API_KEY)
        except Exception as e:
            st.error(f"Failed to initialize Groq client: {e}")
            return False
    
    return True

# --------------------------
# Retry decorator for Qdrant operations
# --------------------------
def retry_qdrant_operation(max_retries=3, delay=2):
    """Retry decorator for Qdrant operations with exponential backoff"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            last_error = None
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    last_error = e
                    error_str = str(e).lower()
                    # Check if it's a retryable error
                    is_retryable = (
                        "timeout" in error_str or 
                        "handshake" in error_str or 
                        "ssl" in error_str or
                        "connection" in error_str
                    )
                    if is_retryable and attempt < max_retries - 1:
                        time.sleep(delay * (attempt + 1))  # Exponential backoff
                        continue
                    else:
                        raise
            raise last_error
        return wrapper
    return decorator

# --------------------------
# Ensure collection exists
# --------------------------
def ensure_collection_exists(silent=False):
    """Create collection if it doesn't exist"""
    if not st.session_state.qdrant:
        return False
    
    try:
        # Check if collection exists with timeout handling
        try:
            st.session_state.qdrant.get_collection(COLLECTION_NAME)
            # Collection exists, return True
            return True
        except Exception as check_error:
            # Collection doesn't exist or check failed
            error_str = str(check_error).lower()
            is_timeout = "timeout" in error_str or "handshake" in error_str or "ssl" in error_str
            
            # For timeout errors, assume collection exists and proceed silently
            if is_timeout:
                return True
            
            # Try to create collection
            try:
                vector_size = get_vector_size()
                st.session_state.qdrant.create_collection(
                    collection_name=COLLECTION_NAME,
                    vectors_config=VectorParams(
                        size=vector_size,
                        distance=Distance.COSINE
                    )
                )
                return True
            except Exception as create_error:
                create_error_str = str(create_error).lower()
                # If collection already exists, that's fine
                if "already exists" in create_error_str or ("collection" in create_error_str and "exists" in create_error_str):
                    return True
                # For timeout errors, assume it exists
                if "timeout" in create_error_str or "handshake" in create_error_str or "ssl" in create_error_str:
                    return True
                # For other errors, only show if not silent
                if not silent:
                    st.error(f"Error ensuring collection exists: {str(create_error)}")
                # Try to proceed anyway
                return True
    except Exception as e:
        # Check if it's a timeout error
        error_str = str(e).lower()
        is_timeout = "timeout" in error_str or "handshake" in error_str or "ssl" in error_str
        
        # For timeout errors, assume collection exists and proceed silently
        if is_timeout:
            return True
        
        # Only show non-timeout errors
        if not silent:
            st.error(f"Error ensuring collection exists: {str(e)}")
        # Return True to allow upload to proceed
        return True

# --------------------------
# SharePoint OAuth2 Authentication
# --------------------------
def get_sharepoint_access_token():
    """Get access token using MSAL device code flow for SharePoint"""
    # Use SharePoint-specific scopes and endpoint
    # For SharePoint Online, we need to use the SharePoint-specific resource
    SCOPES = ["https://graph.microsoft.com/Files.Read.All", "https://graph.microsoft.com/Sites.Read.All"]
    
    # Use a different approach - try with user's tenant or common endpoint
    # For SharePoint REST API, we might need different scopes
    CLIENT_ID = "14d82eec-204b-4c2f-b7e8-296a70dab67e"  # Microsoft Graph Explorer client ID (has broader permissions)
    
    # Alternative: Use SharePoint REST API scopes
    sharepoint_scopes = ["https://[tenant].sharepoint.com/.default"]  # This won't work without tenant
    
    # Try with common endpoint and Graph API scopes
    app = msal.PublicClientApplication(
        CLIENT_ID,
        authority="https://login.microsoftonline.com/common"
    )
    
    # Try to get token from cache first
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            return result["access_token"]
    
    # Use device code flow
    flow = app.initiate_device_flow(scopes=SCOPES)
    if "user_code" not in flow:
        raise Exception("Failed to create device flow")
    
    return flow  # Return flow for UI display

def download_sharepoint_file_with_oauth(url, access_token):
    """Download SharePoint file using OAuth2 access token"""
    headers = {
        'Authorization': f'Bearer {access_token}',
        'Accept': '*/*',
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
    }
    
    # Try to convert SharePoint URL to Microsoft Graph API format
    # SharePoint URLs can be accessed via Graph API
    parsed = urlparse(url)
    
    # Extract site and file path from SharePoint URL
    # Format: https://[tenant].sharepoint.com/sites/[sitename]/...
    if '/sites/' in parsed.path:
        path_parts = parsed.path.split('/sites/')
        if len(path_parts) > 1:
            site_and_path = path_parts[1]
            site_name = site_and_path.split('/')[0]
            file_path = '/'.join(site_and_path.split('/')[1:])
            
            # Try Graph API endpoint
            tenant = parsed.netloc.split('.')[0] if '.' in parsed.netloc else None
            if tenant:
                # Use Graph API to get file download URL
                graph_url = f"https://graph.microsoft.com/v1.0/sites/{parsed.netloc}:/sites/{site_name}:/drive/root:/{file_path}"
                try:
                    file_info = requests.get(graph_url, headers=headers, timeout=30)
                    if file_info.status_code == 200:
                        download_url = file_info.json().get('@microsoft.graph.downloadUrl')
                        if download_url:
                            detected_ext = os.path.splitext(file_path)[1]
                            response = requests.get(download_url, headers=headers, stream=True, timeout=30)
                            response.raise_for_status()
                            return response, detected_ext
                except:
                    pass
    
    # Fallback: Try direct download with token
    download_url, detected_ext = convert_sharepoint_url(url)
    response = requests.get(download_url, headers=headers, stream=True, timeout=30, allow_redirects=True)
    response.raise_for_status()
    
    return response, detected_ext

# --------------------------
# Download file from URL
# --------------------------
def convert_sharepoint_url(url):
    """Convert SharePoint viewer URL to direct download URL"""
    if 'sharepoint.com' in url.lower():
        try:
            from urllib.parse import urlparse, parse_qs, unquote, quote
            parsed = urlparse(url)
            params = parse_qs(parsed.query)
            
            # Extract file name from URL parameters
            file_param = params.get('file', [None])[0]
            file_ext = None
            file_name = None
            
            if file_param:
                file_name = unquote(file_param)
                file_ext = os.path.splitext(file_name)[1]
            
            # Try multiple SharePoint download URL formats
            # Format 1: Direct download using download.aspx
            if '/_layouts/15/Doc.aspx' in parsed.path or '/Doc.aspx' in parsed.path:
                # Get the site path (everything before /_layouts)
                path_parts = parsed.path.split('/_layouts')
                site_path = path_parts[0] if path_parts else ''
                
                # Construct download URL
                # SharePoint often uses this format for direct downloads
                download_url = f"{parsed.scheme}://{parsed.netloc}{site_path}/_layouts/15/download.aspx"
                
                # Add SourceUrl parameter
                source_url_encoded = quote(url, safe='')
                download_url = f"{download_url}?SourceUrl={source_url_encoded}"
                
                return download_url, file_ext
            
            # Format 2: Try webdav format (sometimes works for SharePoint)
            if file_name:
                # Try to construct a direct file path
                path_parts = parsed.path.split('/_layouts')
                site_path = path_parts[0] if path_parts else ''
                file_path_encoded = quote(file_name, safe='')
                download_url = f"{parsed.scheme}://{parsed.netloc}{site_path}/{file_path_encoded}"
                return download_url, file_ext
                
        except Exception as e:
            # If conversion fails, return original URL
            pass
    
    return url, None

def download_file_from_url(url, username=None, password=None):
    """Download a file from URL and return the file path and extension"""
    try:
        # Check if it's a SharePoint URL and convert it
        download_url, detected_ext = convert_sharepoint_url(url)
        is_sharepoint = 'sharepoint.com' in url.lower()
        
        # Send GET request
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Accept': '*/*'
        }
        
        # Create session for SharePoint authentication
        session = requests.Session()
        
        # If SharePoint and credentials provided, try authentication
        if is_sharepoint and username and password:
            # Try multiple authentication methods
            auth_success = False
            response = None
            
            # Method 1: Try NTLM authentication (for on-premises SharePoint)
            try:
                auth = HttpNtlmAuth(username, password)
                # First test with a small request to check authentication
                test_response = session.get(download_url, headers=headers, auth=auth, stream=False, timeout=30, allow_redirects=True)
                # Check if we got HTML (login page) instead of the file
                if test_response.status_code == 200:
                    content_preview = test_response.content[:1000].decode('utf-8', errors='ignore').lower()
                    if '<html' in content_preview or '<form' in content_preview or 'sign in' in content_preview or 'login' in content_preview:
                        auth_success = False
                    else:
                        auth_success = True
                        # Now get the full file with streaming
                        response = session.get(download_url, headers=headers, auth=auth, stream=True, timeout=30, allow_redirects=True)
            except Exception:
                auth_success = False
            
            # Method 2: If NTLM failed, try Basic Auth
            if not auth_success:
                try:
                    auth = HTTPBasicAuth(username, password)
                    test_response = session.get(download_url, headers=headers, auth=auth, stream=False, timeout=30, allow_redirects=True)
                    if test_response.status_code == 200:
                        content_preview = test_response.content[:1000].decode('utf-8', errors='ignore').lower()
                        if '<html' in content_preview or '<form' in content_preview or 'sign in' in content_preview or 'login' in content_preview:
                            raise Exception("Authentication failed. SharePoint may require OAuth2 authentication. Please try getting a shareable link or download the file manually.")
                        auth_success = True
                        # Now get the full file with streaming
                        response = session.get(download_url, headers=headers, auth=auth, stream=True, timeout=30, allow_redirects=True)
                except Exception as e:
                    if "Authentication failed" in str(e):
                        raise
                    # Continue to try without auth
                    pass
            
            # If auth failed, try without credentials (in case user is already logged in via browser)
            if not auth_success:
                response = session.get(download_url, headers=headers, stream=True, timeout=30, allow_redirects=True)
        else:
            response = session.get(download_url, headers=headers, stream=True, timeout=30, allow_redirects=True)
        
        response.raise_for_status()
        
        # Get file extension from URL or Content-Type
        parsed_url = urlparse(response.url)  # Use final URL after redirects
        file_extension = os.path.splitext(parsed_url.path)[1]
        
        # If we detected extension from SharePoint URL, use it
        if detected_ext and not file_extension:
            file_extension = detected_ext
        
        # If still no extension, try to get from Content-Type
        if not file_extension or file_extension == '.aspx':
            content_type = response.headers.get('Content-Type', '')
            content_type_map = {
                'application/pdf': '.pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                'application/msword': '.doc',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                'application/vnd.ms-powerpoint': '.ppt',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                'text/plain': '.txt'
            }
            detected_from_content = content_type_map.get(content_type.split(';')[0])
            if detected_from_content:
                file_extension = detected_from_content
            elif detected_ext:
                file_extension = detected_ext
            else:
                # Try to extract from Content-Disposition header
                content_disposition = response.headers.get('Content-Disposition', '')
                if 'filename=' in content_disposition:
                    import re
                    filename_match = re.search(r'filename[^;=\n]*=(([\'"]).*?\2|[^\s;]+)', content_disposition)
                    if filename_match:
                        filename = filename_match.group(1).strip('\'"')
                        file_extension = os.path.splitext(filename)[1]
        
        # If still no valid extension, default based on detected or PDF
        if not file_extension or file_extension == '.aspx':
            if detected_ext:
                file_extension = detected_ext
            else:
                raise Exception("Could not determine file type from URL. Please ensure the URL points directly to a document file (PDF, Word, PowerPoint, or Text).")
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            # Download in chunks
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    tmp_file.write(chunk)
            temp_path = tmp_file.name
        
        # Validate downloaded file
        file_size = os.path.getsize(temp_path)
        if file_size == 0:
            os.unlink(temp_path)
            raise Exception("Downloaded file is empty. The URL may require authentication or the file may not be accessible.")
        
        # Check if file is actually HTML/ASPX (common with SharePoint authentication pages)
        if file_size < 1024:  # Check first 1KB for HTML content
            with open(temp_path, 'rb') as f:
                first_bytes = f.read(1024)
                if b'<html' in first_bytes.lower() or b'<asp' in first_bytes.lower() or b'<!doctype' in first_bytes.lower():
                    os.unlink(temp_path)
                    raise Exception("Downloaded content appears to be a web page (HTML/ASPX) rather than a document file. This usually means the URL requires authentication. Please try: 1) Getting a direct download link from SharePoint, or 2) Downloading the file manually and uploading from your computer.")
        
        # Validate PPTX files (should start with PK - ZIP signature)
        if file_extension.lower() == '.pptx':
            with open(temp_path, 'rb') as f:
                first_bytes = f.read(4)
                if first_bytes[:2] != b'PK':  # PPTX files are ZIP archives starting with PK
                    os.unlink(temp_path)
                    raise Exception("Downloaded file does not appear to be a valid PowerPoint file. The file may be corrupted or the URL may not point to the actual file.")
        
        return temp_path, file_extension
    except requests.exceptions.RequestException as e:
        raise Exception(f"Failed to download file from URL: {str(e)}. Please ensure the URL is accessible and points to a downloadable file.")
    except Exception as e:
        raise Exception(f"Error processing URL: {str(e)}")

# --------------------------
# Text Extraction from Different File Types
# --------------------------
def extract_text_from_file(file_path, file_extension):
    """Extract text from various file formats"""
    file_extension = file_extension.lower()
    
    if file_extension == '.pdf':
        # Extract text from PDF
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            pages.append(text)
        return "\n".join(pages), len(reader.pages)
    
    elif file_extension in ['.docx', '.doc']:
        # Extract text from Word document
        try:
            doc = DocxDocument(file_path)
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            return "\n".join(paragraphs), len(doc.paragraphs)
        except Exception as e:
            # Try to handle old .doc format (requires additional library)
            raise Exception(f"Error reading Word document. For .doc files, please convert to .docx format. Error: {str(e)}")
    
    elif file_extension in ['.pptx', '.ppt']:
        # Extract text from PowerPoint presentation
        try:
            # Validate file exists and is readable
            if not os.path.exists(file_path):
                raise Exception(f"File not found at path: {file_path}")
            
            # Check file size
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                raise Exception("PowerPoint file is empty or corrupted")
            
            # Try to open the presentation
            prs = Presentation(file_path)
            text_runs = []
            slide_count = len(prs.slides)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_runs.append(shape.text)
                    # Also check for text in tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_runs.append(cell.text)
            
            if not text_runs:
                return "", slide_count
            
            return "\n".join(text_runs), slide_count
        except Exception as e:
            error_msg = str(e)
            if "Package not found" in error_msg or "not a zip file" in error_msg.lower():
                raise Exception(f"Invalid PowerPoint file. The file may be corrupted, incomplete, or not actually a PowerPoint file. If downloading from SharePoint, ensure you have a direct download link. Error: {error_msg}")
            elif file_extension == '.ppt':
                raise Exception(f"Error reading PowerPoint file. Legacy .ppt format is not fully supported. Please convert to .pptx format. Error: {error_msg}")
            else:
                raise Exception(f"Error reading PowerPoint file: {error_msg}")
    
    elif file_extension == '.txt':
        # Read text file directly
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        lines = content.split('\n')
        return content, len([l for l in lines if l.strip()])
    
    else:
        raise Exception(f"Unsupported file format: {file_extension}")

# --------------------------
# Groq LLM
# --------------------------
def generate_llm_answer(query, context):
    prompt = f"""
You are an expert assistant. Use ONLY the context below to answer.

USER QUERY:
{query}

CONTEXT:
{context}

RULES:
- Only use the provided context.
- If context is insufficient, respond: "Not enough information in the document."
"""
    try:
        stream = st.session_state.groq_client.chat.completions.create(
            model=LLM_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_completion_tokens=512,
            stream=True
        )

        answer = ""
        for chunk in stream:
            answer += chunk.choices[0].delta.content or ""
        return answer
    except Exception as e:
        return f"Error generating answer: {str(e)}"

# --------------------------
# Main UI
# --------------------------
st.markdown('<h1 class="main-header">📚 Delivery Management System</h1>', unsafe_allow_html=True)

# Initialize clients
if not init_clients():
    st.stop()

# Sidebar
with st.sidebar:
    st.header("⚙️ Configuration")
    st.info(f"**Collection:** {COLLECTION_NAME}")
    st.info(f"**Embedding Model:** {EMBEDDING_MODEL}")
    st.info(f"**LLM Model:** {LLM_MODEL}")
    
    st.divider()
    
    if st.button("🗑️ Delete Collection", type="secondary", use_container_width=True):
        if st.session_state.qdrant:
            try:
                st.session_state.qdrant.delete_collection(collection_name=COLLECTION_NAME)
                st.success("Collection deleted successfully!")
                st.rerun()
            except Exception as e:
                st.error(f"Error deleting collection: {e}")

# Main content area
tab1, tab2 = st.tabs(["🔍 Search", "📤 Upload Docs"])

with tab1:
    st.header("Ask Questions About Your Documents")
    
    query = st.text_input(
        "Enter your question:",
        placeholder="What information are you looking for?",
        key="search_query"
    )
    
    col1, col2 = st.columns([1, 4])
    with col1:
        search_button = st.button("🔍 Search", type="primary", use_container_width=True)
    
    if search_button and query:
        if not query.strip():
            st.warning("Please enter a question!")
        else:
            with st.spinner("Searching and generating answer..."):
                try:
                    # Ensure collection exists
                    if not ensure_collection_exists():
                        st.warning("Collection doesn't exist. Please upload a document first.")
                    else:
                        # Generate embedding
                        qvec = embed(query).tolist()
                        
                        # Search in Qdrant
                        hits = st.session_state.qdrant.query_points(
                            collection_name=COLLECTION_NAME,
                            query=qvec,
                            limit=MAX_SEARCH_RESULTS,
                            with_payload=True
                        )
                        
                        if not hits.points:
                            st.warning("No results found. Please upload a document first.")
                        else:
                            # Prepare context
                            chunks = [p.payload.get("text", "") for p in hits.points]
                            context = "\n\n".join(chunks)
                            
                            # Generate answer
                            answer = generate_llm_answer(query, context)
                            
                            # Display AI response
                            st.markdown("### 🤖 AI Response")
                            st.markdown(f'<div class="ai-response">{answer}</div>', unsafe_allow_html=True)
                
                except Exception as e:
                    st.error(f"Error during search: {str(e)}")

with tab2:
    st.header("Upload and Ingest Documents")
    
    # Toggle between file upload and URL upload
    upload_method = st.radio(
        "Choose upload method:",
        ["📁 Upload from Computer", "🔗 Upload from URL"],
        horizontal=True
    )
    
    st.divider()
    
    # Chunk size and overlap are set to defaults (hidden from UI)
    chunk_size = 500  # Default chunk size
    overlap = 100     # Default overlap
    
    file_to_process = None
    file_source = None
    filename = None
    
    if upload_method == "📁 Upload from Computer":
        uploaded_file = st.file_uploader(
            "Choose a document file",
            type=["pdf", "docx", "doc", "pptx", "ppt", "txt"],
            help="Upload a document (PDF, Word, PowerPoint, or Text) to add it to the vector database"
        )
        if uploaded_file is not None:
            file_to_process = uploaded_file
            file_source = "upload"
            filename = uploaded_file.name
    else:
        # URL upload
        url_input = st.text_input(
            "Enter document URL:",
            placeholder="https://example.com/document.pdf",
            help="Enter a direct link to a document (PDF, Word, PowerPoint, or Text file)"
        )
        
        # Check if it's a SharePoint URL
        is_sharepoint_url = url_input and 'sharepoint.com' in url_input.lower()
        
        sharepoint_username = None
        sharepoint_password = None
        
        if is_sharepoint_url:
            st.info("🔐 **SharePoint URL detected** - SSO Authentication required")
            
            # Check if we have a token in session
            if 'sharepoint_token' not in st.session_state:
                st.session_state.sharepoint_token = None
            
            auth_method = st.radio(
                "Authentication Method:",
                ["🔐 SSO (OAuth2 - Recommended)", "🔑 Username/Password"],
                horizontal=True
            )
            
            sharepoint_username = None
            sharepoint_password = None
            
            if auth_method == "🔐 SSO (OAuth2 - Recommended)":
                if st.session_state.sharepoint_token:
                    st.success("✅ Authenticated with SSO")
                    if st.button("🔄 Re-authenticate"):
                        st.session_state.sharepoint_token = None
                        st.rerun()
                else:
                    if st.button("🔐 Authenticate with SSO", type="primary"):
                        with st.spinner("Initiating SSO authentication..."):
                            try:
                                flow = get_sharepoint_access_token()
                                
                                # Display device code
                                st.warning("""
                                **⚠️ Admin Consent Required**
                                
                                Your organization requires admin approval for Microsoft Graph API access.
                                
                                **Options:**
                                1. **Contact IT Admin**: Ask them to approve "Microsoft Graph Command Line Tools" app
                                2. **Use Alternative Method**: Switch to "Username/Password" authentication below
                                3. **Get Shareable Link**: In SharePoint, get a "Anyone with the link" shareable URL
                                4. **Download Manually**: Download the file from SharePoint and upload from computer
                                """)
                                
                                st.info(f"""
                                **If you have admin approval, continue:**
                                1. Go to: **{flow['verification_uri']}**
                                2. Enter this code: **{flow['user_code']}**
                                3. Sign in with your Microsoft account
                                4. Wait for confirmation below
                                """)
                                
                                # Poll for token
                                app = msal.PublicClientApplication(
                                    "14d82eec-204b-4c2f-b7e8-296a70dab67e",
                                    authority="https://login.microsoftonline.com/common"
                                )
                                
                                result = None
                                expiration = flow.get("expires_in", 900) + time.time()
                                while time.time() < expiration:
                                    result = app.acquire_token_by_device_flow(flow)
                                    if "access_token" in result:
                                        st.session_state.sharepoint_token = result["access_token"]
                                        st.success("✅ Authentication successful!")
                                        st.rerun()
                                        break
                                    elif "error" in result:
                                        st.error(f"Authentication error: {result.get('error_description', result.get('error'))}")
                                        break
                                    time.sleep(5)
                                
                                if not result or "access_token" not in result:
                                    st.error("Authentication timed out. Please try again.")
                            except Exception as e:
                                error_msg = str(e)
                                if "AADSTS65002" in error_msg or "admin" in error_msg.lower() or "consent" in error_msg.lower():
                                    st.error("""
                                    **Admin Consent Required**
                                    
                                    Your organization requires IT admin approval for this application.
                                    
                                    **Recommended Solutions:**
                                    1. **Contact IT Admin**: Request approval for Microsoft Graph API access
                                    2. **Use Username/Password**: Try the alternative authentication method
                                    3. **Use Shareable Link**: Get a public shareable link from SharePoint
                                    4. **Manual Upload**: Download file from SharePoint and upload from computer
                                    """)
                                else:
                                    st.error(f"SSO authentication failed: {error_msg}")
            else:
                with st.expander("🔑 SharePoint Credentials", expanded=True):
                    sharepoint_username = st.text_input(
                        "Username/Email:",
                        placeholder="your.email@company.com",
                        help="Your SharePoint/Office 365 username"
                    )
                    sharepoint_password = st.text_input(
                        "Password:",
                        type="password",
                        placeholder="Enter your password",
                        help="Your SharePoint/Office 365 password"
                    )
                    st.caption("💡 **Note:** For SharePoint Online, SSO is recommended. Username/Password may not work for modern SharePoint.")
        
        if url_input and url_input.strip():
            file_to_process = url_input.strip()
            file_source = "url"
            filename = os.path.basename(urlparse(url_input).path) or "document_from_url"
    
    if file_to_process:
        if st.button("📤 Upload & Ingest", type="primary", use_container_width=True):
            with st.spinner("Processing document..."):
                try:
                    temp_path = None
                    
                    if file_source == "url":
                        # Download file from URL
                        with st.spinner("Downloading file from URL..."):
                            # Check if SharePoint with SSO token
                            if is_sharepoint_url and 'sharepoint_token' in st.session_state and st.session_state.sharepoint_token:
                                try:
                                    response, detected_ext = download_sharepoint_file_with_oauth(file_to_process, st.session_state.sharepoint_token)
                                    # Save to temp file
                                    file_extension = detected_ext or '.pptx'
                                    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                                        for chunk in response.iter_content(chunk_size=8192):
                                            if chunk:
                                                tmp_file.write(chunk)
                                        temp_path = tmp_file.name
                                except Exception as e:
                                    if "401" in str(e) or "Unauthorized" in str(e) or "403" in str(e):
                                        st.session_state.sharepoint_token = None
                                        st.error("Authentication expired or insufficient permissions. Please re-authenticate with SSO.")
                                        st.stop()
                                    else:
                                        raise
                            else:
                                # Pass credentials if SharePoint and provided
                                username = sharepoint_username if is_sharepoint_url and sharepoint_username else None
                                password = sharepoint_password if is_sharepoint_url and sharepoint_password else None
                                temp_path, file_extension = download_file_from_url(file_to_process, username, password)
                    else:
                        # Process uploaded file
                        file_extension = os.path.splitext(file_to_process.name)[1]
                        
                        # Save uploaded file temporarily with correct extension
                        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                            tmp_file.write(file_to_process.read())
                            temp_path = tmp_file.name
                    
                    # Determine file type name
                    file_type_name = {
                        '.pdf': 'PDF',
                        '.docx': 'Word Document',
                        '.doc': 'Word Document (Legacy)',
                        '.pptx': 'PowerPoint Presentation',
                        '.ppt': 'PowerPoint Presentation (Legacy)',
                        '.txt': 'Text File'
                    }.get(file_extension.lower(), 'Document')
                    
                    # Extract text from file based on type
                    full_text, page_count = extract_text_from_file(temp_path, file_extension)
                    
                    if not full_text.strip():
                        st.error(f"No text could be extracted from the {file_type_name}.")
                    else:
                        # Chunk the text
                        chunks = []
                        start = 0
                        
                        while start < len(full_text):
                            end = start + chunk_size
                            chunks.append(full_text[start:end])
                            start += chunk_size - overlap
                        
                        # Generate embeddings and store
                        progress_bar = st.progress(0)
                        points = []
                        
                        for i, chunk in enumerate(chunks):
                            vec = embed(chunk).tolist()
                            points.append({
                                "id": int(np.random.randint(1_000_000_000)),
                                "vector": vec,
                                "payload": {"text": chunk}
                            })
                            progress_bar.progress((i + 1) / len(chunks))
                        
                        # Ensure collection exists before upserting (silent mode to avoid showing timeout errors)
                        ensure_collection_exists(silent=True)
                        
                        # Try to upsert with retry logic for timeout errors
                        max_retries = 3
                        upsert_success = False
                        last_error = None
                        
                        for attempt in range(max_retries):
                            try:
                                st.session_state.qdrant.upsert(
                                    collection_name=COLLECTION_NAME,
                                    points=points
                                )
                                upsert_success = True
                                break
                            except Exception as upsert_error:
                                last_error = upsert_error
                                error_str = str(upsert_error).lower()
                                
                                # Check if it's a timeout/SSL error (retryable)
                                is_timeout = "timeout" in error_str or "handshake" in error_str or "ssl" in error_str
                                
                                # Check if collection doesn't exist
                                is_not_found = "not found" in error_str or "doesn't exist" in error_str or "404" in error_str
                                
                                if is_not_found:
                                    # Collection doesn't exist, create it
                                    vector_size = get_vector_size()
                                    try:
                                        st.session_state.qdrant.create_collection(
                                            collection_name=COLLECTION_NAME,
                                            vectors_config=VectorParams(
                                                size=vector_size,
                                                distance=Distance.COSINE
                                            )
                                        )
                                        # Retry upsert after creating collection
                                        continue
                                    except Exception as create_err:
                                        create_err_str = str(create_err).lower()
                                        # If collection already exists, try upsert again
                                        if "already exists" in create_err_str:
                                            continue
                                        # If timeout during creation, retry
                                        if "timeout" in create_err_str or "handshake" in create_err_str:
                                            if attempt < max_retries - 1:
                                                time.sleep(2 * (attempt + 1))
                                                continue
                                
                                # For timeout errors, retry with exponential backoff
                                if is_timeout and attempt < max_retries - 1:
                                    time.sleep(2 * (attempt + 1))  # 2s, 4s, 6s
                                    continue
                                
                                # For other errors or last attempt, raise
                                if attempt == max_retries - 1:
                                    raise
                        
                        if not upsert_success:
                            raise last_error
                        
                        # Clean up temp file
                        os.unlink(temp_path)
                        
                        st.success(f"✅ {file_type_name} successfully ingested!")
                        st.info(f"**Filename:** {filename}")
                        st.info(f"**File Type:** {file_type_name}")
                        if file_source == "url":
                            st.info(f"**Source:** URL")
                        if file_extension in ['.pdf', '.pptx', '.ppt']:
                            st.info(f"**Pages/Slides:** {page_count}")
                        elif file_extension in ['.docx', '.doc']:
                            st.info(f"**Paragraphs:** {page_count}")
                        elif file_extension == '.txt':
                            st.info(f"**Lines:** {page_count}")
                        st.info(f"**Chunks created:** {len(chunks)}")
                        st.info(f"**Total characters:** {len(full_text):,}")
                
                except Exception as e:
                    st.error(f"Error processing document: {str(e)}")
                    if 'temp_path' in locals() and temp_path:
                        try:
                            os.unlink(temp_path)
                        except:
                            pass

# Footer
st.divider()
st.markdown(
    "<div style='text-align: center; color: #666; padding: 1rem;'>"
    "Built with Streamlit | Powered by Qdrant & Groq"
    "</div>",
    unsafe_allow_html=True
)
